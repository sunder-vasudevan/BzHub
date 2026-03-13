#!/usr/bin/env python3
"""
Generates a McKinsey/PwC-style 10-slide executive deck for the BzHub Efficiency Case Study.
Run:   python3.9 documentation/generate_exec_deck.py
Output: documentation/BzHub_Efficiency_ExecDeck.pptx
Requires: pip3 install python-pptx
"""

import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import os

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_NAVY  = RGBColor(0x00, 0x2B, 0x5C)
MID_BLUE   = RGBColor(0x00, 0x5B, 0x99)
TEAL       = RGBColor(0x00, 0x7A, 0x87)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
MID_GREY   = RGBColor(0x88, 0x88, 0x88)
AMBER      = RGBColor(0xE8, 0x8C, 0x00)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)


# ── Low-level helpers ─────────────────────────────────────────────────────────

def blank_slide(prs):
    """Add a completely blank slide."""
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


def rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a filled rectangle shape."""
    from pptx.util import Pt as _Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, left, top, width, height,
          text='', size=18, bold=False, italic=False,
          color=DARK_TEXT, align=PP_ALIGN.LEFT,
          wrap=True, line_spacing=None):
    """Add a text box with single-run text."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        from pptx.util import Pt as _Pt
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return box, tf


def add_logo(slide, path, left, top, height):
    """Add image if file exists, silently skip if not."""
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, height=height)


def footer_bar(slide, label='BzHub Efficiency Case Study  |  March 2026  |  Confidential'):
    """Thin navy footer strip at bottom of every slide."""
    rect(slide, 0, H - Inches(0.3), W, Inches(0.3), fill_color=DARK_NAVY)
    txbox(slide, Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
          text=label, size=Pt(7), color=WHITE, align=PP_ALIGN.LEFT)


def section_label(slide, text, left=Inches(0.5), top=Inches(0.18)):
    """Small coloured section tag at top-left."""
    rect(slide, left, top, Inches(2.8), Inches(0.22), fill_color=TEAL)
    txbox(slide, left + Inches(0.08), top, Inches(2.6), Inches(0.22),
          text=text.upper(), size=Pt(7.5), bold=True, color=WHITE)


def slide_title(slide, text, top=Inches(0.52), color=DARK_NAVY):
    txbox(slide, Inches(0.5), top, Inches(12.3), Inches(0.6),
          text=text, size=Pt(26), bold=True, color=color)


def rule_line(slide, top, color=MID_BLUE, left=Inches(0.5), width=None):
    """Thin horizontal coloured rule."""
    w = width or (W - Inches(1.0))
    r = rect(slide, left, top, w, Inches(0.025), fill_color=color)


def big_number(slide, number, label, left, top, num_color=MID_BLUE):
    """Large KPI number with small label below."""
    txbox(slide, left, top, Inches(2.8), Inches(1.1),
          text=number, size=Pt(52), bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txbox(slide, left, top + Inches(1.0), Inches(2.8), Inches(0.35),
          text=label, size=Pt(10), color=MID_GREY, align=PP_ALIGN.CENTER)


def bullet_block(slide, items, left, top, width, size=Pt(11), color=DARK_TEXT,
                 bullet_char='▸', spacing=Inches(0.32)):
    """Render a list of bullet strings as individual text boxes (reliable spacing)."""
    for i, item in enumerate(items):
        txbox(slide, left, top + i * spacing, width, Inches(0.3),
              text=f'{bullet_char}  {item}', size=size, color=color)


def callout(slide, text, left, top, width, height,
            bg=RGBColor(0xE8, 0xF4, 0xF6), border=TEAL):
    """Teal-bordered insight box."""
    rect(slide, left, top, width, height, fill_color=bg, line_color=border, line_width=Pt(2))
    txbox(slide, left + Inches(0.15), top + Inches(0.12),
          width - Inches(0.3), height - Inches(0.24),
          text=text, size=Pt(11), italic=True, color=DARK_NAVY, wrap=True)


def two_col_table(slide, headers, rows, left, top, col_widths, row_height=Inches(0.32)):
    """Simple table rendered as coloured rectangles + text boxes."""
    n_cols = len(headers)
    # Header row
    x = left
    for i, h in enumerate(headers):
        rect(slide, x, top, col_widths[i], row_height, fill_color=MID_BLUE)
        txbox(slide, x + Inches(0.08), top + Inches(0.05),
              col_widths[i] - Inches(0.1), row_height - Inches(0.08),
              text=h, size=Pt(9), bold=True, color=WHITE)
        x += col_widths[i]
    # Data rows
    for r_idx, row in enumerate(rows):
        bg = LIGHT_GREY if r_idx % 2 == 0 else WHITE
        x = left
        for c_idx, cell in enumerate(row):
            rect(slide, x, top + (r_idx + 1) * row_height,
                 col_widths[c_idx], row_height, fill_color=bg)
            txbox(slide, x + Inches(0.08),
                  top + (r_idx + 1) * row_height + Inches(0.05),
                  col_widths[c_idx] - Inches(0.1), row_height - Inches(0.08),
                  text=str(cell), size=Pt(9), color=DARK_TEXT)
            x += col_widths[c_idx]


# ── Chart helpers ─────────────────────────────────────────────────────────────

def _c(rgb_tuple):
    return tuple(c/255 for c in rgb_tuple)

_NAVY  = (0x00, 0x2B, 0x5C)
_BLUE  = (0x00, 0x5B, 0x99)
_TEAL  = (0x00, 0x7A, 0x87)
_LGREY = (0xF2, 0xF4, 0xF7)
_DTEXT = (0x1A, 0x1A, 0x2E)


def _savefig(fig, dpi=150):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=dpi, bbox_inches='tight', facecolor='none')
    buf.seek(0)
    plt.close(fig)
    return buf


def deck_chart_sessions():
    """Grouped bar: session time (left axis) vs. releases shipped (right axis)."""
    sessions = ['Session 1\n(45 min)', 'Session 2\n(25 min)', 'Session 3\n(35 min)']
    times    = [45, 25, 35]
    releases = [0, 1, 6]
    x = np.arange(len(sessions))
    width = 0.38

    fig, ax1 = plt.subplots(figsize=(5.8, 3.0), facecolor='none')
    ax2 = ax1.twinx()
    ax1.set_facecolor('none')
    fig.patch.set_alpha(0)

    b1 = ax1.bar(x - width/2, times, width, color=_c(_BLUE), edgecolor='white', linewidth=0.8)
    b2 = ax2.bar(x + width/2, releases, width, color=_c(_TEAL), edgecolor='white', linewidth=0.8)

    for bar, val in zip(b1, times):
        ax1.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.8,
                 f'{val}m', ha='center', va='bottom', fontsize=9, fontweight='bold',
                 color=_c(_BLUE))
    for bar, val in zip(b2, releases):
        ax2.text(bar.get_x() + bar.get_width()/2, bar.get_height() + 0.08,
                 str(val), ha='center', va='bottom', fontsize=9, fontweight='bold',
                 color=_c(_TEAL))

    ax1.set_ylabel('Interaction Time (mins)', fontsize=8, color=_c(_BLUE))
    ax2.set_ylabel('Releases Shipped', fontsize=8, color=_c(_TEAL))
    ax1.set_title('Session Efficiency', fontsize=10, fontweight='bold',
                  color=_c(_NAVY), pad=8)
    ax1.set_xticks(x)
    ax1.set_xticklabels(sessions, fontsize=8)
    ax1.set_ylim(0, 60)
    ax2.set_ylim(0, 8)
    for ax in [ax1, ax2]:
        ax.spines['top'].set_visible(False)
    ax1.spines['left'].set_color('#AAAAAA')
    ax1.spines['bottom'].set_color('#AAAAAA')
    ax1.tick_params(labelsize=8)
    ax2.tick_params(labelsize=8)

    lines1, labels1 = ax1.get_legend_handles_labels()
    b1_patch = mpatches.Patch(color=_c(_BLUE), label='Interaction Time (mins)')
    b2_patch = mpatches.Patch(color=_c(_TEAL), label='Releases Shipped')
    ax1.legend(handles=[b1_patch, b2_patch], fontsize=7, loc='upper left',
               framealpha=0.4, facecolor='white')
    fig.tight_layout()
    return _savefig(fig)


def add_deck_picture(slide, buf, left, top, width):
    """Insert a BytesIO chart image into a pptx slide."""
    slide.shapes.add_picture(buf, left, top, width=width)


LOGO = '/Users/scottvalentino/BzHub/assets/bizhub_logo.png'


# ── SLIDE BUILDERS ────────────────────────────────────────────────────────────

def slide_01_cover(prs):
    """Cover — dark navy full bleed."""
    s = blank_slide(prs)
    # Full navy background
    rect(s, 0, 0, W, H, fill_color=DARK_NAVY)
    # Accent bar (teal, left edge)
    rect(s, 0, 0, Inches(0.18), H, fill_color=TEAL)
    # Top rule
    rect(s, Inches(0.5), Inches(1.6), W - Inches(1.0), Inches(0.04), fill_color=TEAL)

    # Title
    txbox(s, Inches(0.6), Inches(1.8), Inches(11.5), Inches(1.0),
          text='Engineering Efficiency with AI Pair Programming',
          size=Pt(34), bold=True, color=WHITE)
    # Subtitle
    txbox(s, Inches(0.6), Inches(2.9), Inches(10.0), Inches(0.5),
          text='A Case Study of BzHub — One Product Owner, One Claude',
          size=Pt(18), color=RGBColor(0xA8, 0xC8, 0xE8))
    # Bottom rule
    rect(s, Inches(0.5), Inches(4.0), Inches(4.0), Inches(0.04), fill_color=MID_BLUE)
    # Meta
    meta = [
        ('Author', 'sunder-vasudevan'),
        ('Date', 'March 2026'),
        ('Version', '1.1 — with measured interaction time data'),
        ('Classification', 'Technical Peers'),
    ]
    for i, (k, v) in enumerate(meta):
        txbox(s, Inches(0.6), Inches(4.2) + i * Inches(0.38),
              Inches(1.6), Inches(0.32),
              text=k.upper(), size=Pt(8), bold=True, color=TEAL)
        txbox(s, Inches(2.3), Inches(4.2) + i * Inches(0.38),
              Inches(5.0), Inches(0.32),
              text=v, size=Pt(8), color=RGBColor(0xCC, 0xDD, 0xEE))
    # Logo (bottom right)
    add_logo(s, LOGO, W - Inches(2.2), H - Inches(1.6), Inches(1.4))
    # Footer
    rect(s, 0, H - Inches(0.3), W, Inches(0.3), fill_color=RGBColor(0x00, 0x1A, 0x3A))
    txbox(s, Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
          text='© 2026 sunder-vasudevan. All rights reserved.',
          size=Pt(7), color=RGBColor(0x88, 0x99, 0xAA))


def slide_02_context(prs):
    """Slide 2 — The Problem / Context."""
    s = blank_slide(prs)
    section_label(s, '01  Context')
    slide_title(s, 'The Challenge: Building a Full Product Alone Is Hard')
    rule_line(s, Inches(1.2))

    # Left column — the old reality
    rect(s, Inches(0.5), Inches(1.4), Inches(5.6), Inches(4.8), fill_color=LIGHT_GREY)
    txbox(s, Inches(0.7), Inches(1.55), Inches(5.2), Inches(0.35),
          text='TRADITIONAL SOLO DEVELOPER', size=Pt(9), bold=True, color=MID_GREY)
    items_left = [
        'Context-switches kill momentum',
        'Full-stack gaps slow delivery',
        'Documentation is always last',
        'Architecture decisions made alone',
        'Cognitive fatigue limits sprint length',
        'Hiring a team = 3–6 month delay + cost',
    ]
    bullet_block(s, items_left, Inches(0.7), Inches(2.0), Inches(5.0),
                 size=Pt(10.5), color=DARK_TEXT, spacing=Inches(0.42))

    # Right column — the question
    rect(s, Inches(6.5), Inches(1.4), Inches(6.3), Inches(4.8), fill_color=RGBColor(0xE8, 0xF2, 0xFC))
    txbox(s, Inches(6.7), Inches(1.55), Inches(5.9), Inches(0.35),
          text='THE QUESTION THIS STUDY ANSWERS', size=Pt(9), bold=True, color=MID_BLUE)
    callout(s,
            'Can one Product Owner working with Claude as a full-session AI '
            'pair-programmer match the output of a traditional 3-person startup '
            'development team — in velocity, scope, and quality?',
            Inches(6.7), Inches(2.0), Inches(5.8), Inches(1.5))
    txbox(s, Inches(6.7), Inches(3.65), Inches(5.8), Inches(0.35),
          text='SPOILER', size=Pt(9), bold=True, color=AMBER)
    txbox(s, Inches(6.7), Inches(3.95), Inches(5.8), Inches(1.1),
          text='Yes — 2–4x velocity multiplier confirmed by git data and measured logs: 105 mins of interaction → 6 production releases → 60–120x compression ratio.',
          size=Pt(12), bold=True, color=DARK_NAVY, wrap=True)

    footer_bar(s)


def slide_03_experiment(prs):
    """Slide 3 — The Experiment."""
    s = blank_slide(prs)
    section_label(s, '02  The Experiment')
    slide_title(s, 'BzHub: A Real Product, Built in 31 Days')
    rule_line(s, Inches(1.2))

    # KPI row
    kpis = [
        ('31', 'Calendar days'),
        ('~10', 'Active coding days'),
        ('64', 'Total commits'),
        ('11', 'Production releases'),
        ('9', 'Live modules'),
        ('~15,800', 'Net new LOC'),
    ]
    for i, (num, lbl) in enumerate(kpis):
        big_number(s, num, lbl,
                   left=Inches(0.4) + i * Inches(2.1),
                   top=Inches(1.5))

    rule_line(s, Inches(3.1), color=LIGHT_GREY)

    # Architecture evolution
    txbox(s, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.3),
          text='ARCHITECTURE EVOLUTION', size=Pt(9), bold=True, color=MID_GREY)

    stages = [
        ('v1.0\nFeb 10', 'Desktop App\nTkinter / Python'),
        ('v2.0\nFeb 18', 'Web + API\nNext.js + FastAPI\n+ SQLite'),
        ('v4.0\nMar 10', 'Cloud SaaS\nNext.js + Supabase\n+ Vercel'),
    ]
    arrow_y = Inches(3.65)
    for i, (version, desc) in enumerate(stages):
        x = Inches(0.5) + i * Inches(4.1)
        rect(s, x, arrow_y, Inches(3.5), Inches(1.5), fill_color=MID_BLUE)
        txbox(s, x + Inches(0.12), arrow_y + Inches(0.12), Inches(3.26), Inches(0.4),
              text=version, size=Pt(9), bold=True, color=WHITE)
        txbox(s, x + Inches(0.12), arrow_y + Inches(0.5), Inches(3.26), Inches(0.9),
              text=desc, size=Pt(10), color=RGBColor(0xC0, 0xD8, 0xF0), wrap=True)
        if i < 2:
            txbox(s, x + Inches(3.55), arrow_y + Inches(0.55), Inches(0.5), Inches(0.4),
                  text='→', size=Pt(22), bold=True, color=TEAL)

    txbox(s, Inches(0.5), Inches(5.35), Inches(12.3), Inches(0.3),
          text='Each migration was deliberate, documented, and completed in a single sprint without breaking the deployed state.',
          size=Pt(9), italic=True, color=MID_GREY)

    footer_bar(s)


def slide_04_what_was_built(prs):
    """Slide 4 — What Was Built."""
    s = blank_slide(prs)
    section_label(s, '03  Scope')
    slide_title(s, '9 Live Modules. Full Stack. Cloud-Deployed.')
    rule_line(s, Inches(1.2))

    modules = [
        ('Dashboard',             'KPI cards, trend chart, Smart Insights, customizable layout'),
        ('Operations',            'Inventory, POS, Bills, Suppliers, Purchase Orders + approval'),
        ('HR',                    'Employees, Payroll, Goals, Appraisals, Skills, Leave + workflows'),
        ('CRM',                   'Contacts, Leads — List / Kanban / Funnel views, lead scoring'),
        ('Reports',               'Sales Report, Top Sellers, Inventory Report + CSV export'),
        ('Settings',              'Company info, currency, industry templates, brand color, Custom Fields'),
        ('Help',                  'In-app user guide for all modules'),
        ('Employee Self-Service', 'My Goals, My Appraisals, My Leave, My Skills'),
        ('Notifications',         'Bell icon: pending leave / PO / appraisal / low-stock alerts'),
    ]

    for i, (mod, desc) in enumerate(modules):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.45) + row * Inches(1.3)
        rect(s, x, y, Inches(5.9), Inches(1.1),
             fill_color=LIGHT_GREY if col == 0 else RGBColor(0xE8, 0xF2, 0xFC))
        txbox(s, x + Inches(0.15), y + Inches(0.1), Inches(5.6), Inches(0.32),
              text=mod.upper(), size=Pt(9), bold=True, color=MID_BLUE)
        txbox(s, x + Inches(0.15), y + Inches(0.42), Inches(5.6), Inches(0.55),
              text=desc, size=Pt(10), color=DARK_TEXT, wrap=True)

    footer_bar(s)


def slide_05_velocity(prs):
    """Slide 5 — The 2-Day Sprint (velocity evidence)."""
    s = blank_slide(prs)
    section_label(s, '04  Velocity')
    slide_title(s, 'The 2-Day Sprint: 10 Releases, ~7,200 Lines')
    rule_line(s, Inches(1.2))

    # Table (left side — 5 cols, compact)
    two_col_table(s,
        headers=['Release', 'Feature', 'LOC'],
        rows=[
            ('v4.1.x–4.5.0', 'Cloud, Reports, HR, Approvals, Employee Portal (5 releases)', '~4,372'),
            ('v4.6.0',       'Notification Center, CSV Export, Global Search, Audit Log',    '1,267'),
            ('v4.7.0 ●',     'Industry Templates  [Session 2 — 25 min]',                     '169'),
            ('v4.7.1 ●',     'Dynamic brand color — 24 files  [Session 3 — 35 min]',          '971'),
            ('v4.8.0 ●',     'Smart Insights dashboard card  [Session 3]',                   '196'),
            ('v4.9.x ●',     'CRM view switcher: List/Kanban/Funnel  [Session 3]',            '264'),
            ('TOTAL',        '10 releases  ·  all layers  ·  2 calendar days',               '~7,239'),
        ],
        left=Inches(0.5), top=Inches(1.35),
        col_widths=[Inches(1.6), Inches(7.3), Inches(1.1)],
        row_height=Inches(0.38)
    )

    # Callout right
    rect(s, Inches(10.6), Inches(1.35), Inches(2.4), Inches(3.04),
         fill_color=DARK_NAVY)
    txbox(s, Inches(10.7), Inches(1.48), Inches(2.2), Inches(0.4),
          text='IN 2 DAYS', size=Pt(9), bold=True, color=TEAL)
    txbox(s, Inches(10.7), Inches(1.88), Inches(2.2), Inches(0.65),
          text='10', size=Pt(56), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, Inches(10.7), Inches(2.53), Inches(2.2), Inches(0.35),
          text='production releases', size=Pt(10), color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    txbox(s, Inches(10.7), Inches(2.88), Inches(2.2), Inches(0.35),
          text='shipped to cloud', size=Pt(10), color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    txbox(s, Inches(10.7), Inches(3.3), Inches(2.2), Inches(0.35),
          text='● = logged session', size=Pt(8), color=TEAL, align=PP_ALIGN.CENTER)

    # Bottom callout
    callout(s,
            'Sessions 2 + 3 (60 mins combined interaction time) produced 7 of these 10 releases. '
            'A 3-person team would realistically require 3–6 two-week sprints for this scope. '
            'The AI-assisted pair compressed it to 2 calendar days.',
            Inches(0.5), Inches(4.62), Inches(12.3), Inches(0.72))

    footer_bar(s)


def slide_05b_measured_time(prs):
    """Slide 5b — Measured Interaction Time (new: logged session data)."""
    s = blank_slide(prs)
    section_label(s, '04B  Measured Time')
    slide_title(s, '105 Minutes of Interaction → 6 Production Releases')
    rule_line(s, Inches(1.2))

    # 3 session cards
    sessions = [
        ('Session 1', '~45 min', '—', 'Efficiency whitepaper,\ndoc generators'),
        ('Session 2', '~25 min', 'v4.7.0', 'Industry Templates\n(FEAT-038)'),
        ('Session 3', '~35 min', 'v4.7.1 – v4.9.2\n(6 releases)', 'Brand color, Smart Insights,\nCRM List/Kanban/Funnel'),
    ]
    colors = [DARK_NAVY, MID_BLUE, TEAL]
    for i, (title, time, releases, output) in enumerate(sessions):
        x = Inches(0.5) + i * Inches(3.9)
        rect(s, x, Inches(1.4), Inches(3.6), Inches(0.36), fill_color=colors[i])
        txbox(s, x + Inches(0.12), Inches(1.42), Inches(3.38), Inches(0.3),
              text=title, size=Pt(10), bold=True, color=WHITE)
        rect(s, x, Inches(1.76), Inches(3.6), Inches(2.44),
             fill_color=LIGHT_GREY if i < 2 else RGBColor(0xE0, 0xF4, 0xF6))
        txbox(s, x + Inches(0.12), Inches(1.85), Inches(1.2), Inches(0.25),
              text='TIME', size=Pt(7.5), bold=True, color=MID_GREY)
        txbox(s, x + Inches(0.12), Inches(2.1), Inches(3.3), Inches(0.4),
              text=time, size=Pt(22), bold=True, color=colors[i])
        txbox(s, x + Inches(0.12), Inches(2.6), Inches(1.2), Inches(0.25),
              text='RELEASES', size=Pt(7.5), bold=True, color=MID_GREY)
        txbox(s, x + Inches(0.12), Inches(2.85), Inches(3.3), Inches(0.5),
              text=releases, size=Pt(10.5), bold=True, color=DARK_NAVY, wrap=True)
        txbox(s, x + Inches(0.12), Inches(3.42), Inches(1.2), Inches(0.25),
              text='OUTPUT', size=Pt(7.5), bold=True, color=MID_GREY)
        txbox(s, x + Inches(0.12), Inches(3.68), Inches(3.3), Inches(0.45),
              text=output, size=Pt(9.5), color=DARK_TEXT, wrap=True)

    # Total bar
    rect(s, Inches(0.5), Inches(4.28), Inches(11.7), Inches(0.44), fill_color=DARK_NAVY)
    txbox(s, Inches(0.7), Inches(4.33), Inches(11.3), Inches(0.34),
          text='TOTAL:  ~105 mins interaction  →  6 production releases  →  ~1,600 net LOC',
          size=Pt(10.5), bold=True, color=WHITE)

    # Chart (right side)
    chart_buf = deck_chart_sessions()
    add_deck_picture(s, chart_buf, left=Inches(7.8), top=Inches(4.82), width=Inches(5.0))

    # Key stat left
    callout(s,
            'Session 3: 35 mins → 6 releases\n'
            '≈ 3–6 engineer-weeks of traditional work\n'
            '→ 60–120x compression ratio on interaction time',
            Inches(0.5), Inches(4.85), Inches(7.1), Inches(0.96))

    footer_bar(s)


def slide_06_comparison(prs):
    """Slide 6 — Head-to-head comparison."""
    s = blank_slide(prs)
    section_label(s, '05  Comparison')
    slide_title(s, '1 Product Owner + Claude vs. 3-Person Team')
    rule_line(s, Inches(1.2))

    two_col_table(s,
        headers=['Dimension', '1 PO + Claude  (31 days)', '3-Person Team  (31 days est.)'],
        rows=[
            ('Major releases shipped',      '11  (v1.0 → v5.0.0)',          '3–5  (1–2 per 2-week sprint)'),
            ('Architecture migrations',     '3 full paradigm shifts',        'Likely 1  (dedicated project)'),
            ('Live modules delivered',      '9',                             '4–6'),
            ('Documentation',              'Per-feature  (project rule)',    'Varies — often deferred'),
            ('Deployment pipeline',        'Fully automated  (Vercel CI)',   '2–5 days setup for a team'),
            ('Measured PO interaction',    '~105 mins  (3 logged sessions)', 'N/A — team always staffed'),
            ('Handoff overhead',           'Zero',                           '15–25% of sprint capacity'),
        ],
        left=Inches(0.5), top=Inches(1.35),
        col_widths=[Inches(3.0), Inches(4.5), Inches(4.5)],
        row_height=Inches(0.42)
    )

    txbox(s, Inches(0.5), Inches(5.32), Inches(12.0), Inches(0.28),
          text='* 3-person team estimates based on conventional startup sprint norms, not empirical benchmarks.',
          size=Pt(8), italic=True, color=MID_GREY)

    callout(s,
            'Key finding: One Product Owner + Claude approximates the feature throughput of a '
            '3-person startup team. Velocity multiplier: 2–4x at project level. Session 3 alone '
            'demonstrates a 60–120x compression ratio on interaction time.',
            Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.72))

    footer_bar(s)


def slide_07_how_it_works(prs):
    """Slide 7 — The 4 structural gains."""
    s = blank_slide(prs)
    section_label(s, '06  Mechanism')
    slide_title(s, 'Where the Gains Come From: 4 Structural Advantages')
    rule_line(s, Inches(1.2))

    advantages = [
        ('Zero Handoff Overhead',
         'Full-stack features are conceived, designed, and implemented in a single session. '
         'No PRs, Slack threads, or standup latency.'),
        ('Persistent Cross-Session Context',
         'Structured NOTES.md + memory files preserve full project state. Claude resumes '
         'with context a new hire would take 1–2 weeks to acquire.'),
        ('On-Demand Full-Stack Competency',
         'Competent first-draft implementations across Next.js, Supabase, Recharts, and '
         'Tailwind — eliminating the "look this up for 45 minutes" tax.'),
        ('Documentation as a Side-Effect',
         'Release notes, help page, and memory files updated after every feature — '
         'a discipline that is rare solo and hard to sustain in teams under pressure.'),
    ]

    colors = [MID_BLUE, TEAL, DARK_NAVY, RGBColor(0x5B, 0x2D, 0x8E)]

    for i, (title, desc) in enumerate(advantages):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.45) + row * Inches(2.55)
        # Number badge
        rect(s, x, y, Inches(0.48), Inches(0.48), fill_color=colors[i])
        txbox(s, x, y, Inches(0.48), Inches(0.48),
              text=str(i + 1), size=Pt(16), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        # Title
        txbox(s, x + Inches(0.6), y, Inches(5.3), Inches(0.46),
              text=title, size=Pt(12), bold=True, color=colors[i])
        # Rule under title
        rule_line(s, y + Inches(0.5), color=LIGHT_GREY, left=x + Inches(0.6), width=Inches(5.3))
        # Body
        txbox(s, x + Inches(0.6), y + Inches(0.58), Inches(5.3), Inches(1.7),
              text=desc, size=Pt(10.5), color=DARK_TEXT, wrap=True)

    footer_bar(s)


def slide_08_quality(prs):
    """Slide 8 — Code quality & architecture signal."""
    s = blank_slide(prs)
    section_label(s, '07  Quality')
    slide_title(s, 'Speed Did Not Come at the Cost of Quality')
    rule_line(s, Inches(1.2))

    # Left: quality signals
    txbox(s, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.32),
          text='QUALITY SIGNALS IN THE GIT HISTORY', size=Pt(9), bold=True, color=MID_GREY)
    signals = [
        'Each architecture migration preceded by a documented decision commit',
        'Every migration left codebase in stable, deployed state before the next began',
        'Known technical debt explicitly logged — not hidden — with a resolution plan',
        'Per-feature documentation enforced as a hard project rule (5-step checklist)',
        'Intentional debt (hardcoded auth, open RLS) acknowledged with planned fix dates',
    ]
    bullet_block(s, signals, Inches(0.5), Inches(1.85), Inches(5.8),
                 size=Pt(10.5), spacing=Inches(0.5))

    # Right: debt table
    txbox(s, Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.32),
          text='INTENTIONAL DEBT — TRACKED, NOT HIDDEN', size=Pt(9), bold=True, color=MID_GREY)
    two_col_table(s,
        headers=['Known Gap', 'Planned Fix'],
        rows=[
            ('Hardcoded auth\n(admin/admin123)', 'FEAT-036: Supabase Auth\n(Phase 3)'),
            ('Open RLS policies\n(no user filter)', 'Tightens when auth ships'),
            ('No organization_id\non tables', 'FEAT-037: Multi-Tenancy\n(Phase 3)'),
        ],
        left=Inches(7.0), top=Inches(1.85),
        col_widths=[Inches(2.9), Inches(2.9)],
        row_height=Inches(0.6)
    )

    # Bottom callout
    callout(s,
            '"This reflects the discipline of an experienced Product Owner, not the output '
            'of unchecked code generation. Architecture was never sacrificed for speed."',
            Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.72))

    footer_bar(s)


def slide_09_conclusion(prs):
    """Slide 9 — Conclusion + key finding."""
    s = blank_slide(prs)
    # Dark navy background for impact
    rect(s, 0, 0, W, H, fill_color=DARK_NAVY)
    rect(s, 0, 0, Inches(0.18), H, fill_color=TEAL)

    txbox(s, Inches(0.6), Inches(0.35), Inches(3.0), Inches(0.25),
          text='08  CONCLUSION', size=Pt(8), bold=True, color=TEAL)

    txbox(s, Inches(0.6), Inches(0.7), Inches(12.0), Inches(0.9),
          text='The Verdict', size=Pt(32), bold=True, color=WHITE)

    rect(s, Inches(0.6), Inches(1.65), Inches(8.0), Inches(0.04), fill_color=TEAL)

    # Main statement
    txbox(s, Inches(0.6), Inches(1.85), Inches(11.5), Inches(1.3),
          text='One Product Owner + Claude ≈ 3-Person Startup Dev Team',
          size=Pt(24), bold=True, color=RGBColor(0xA8, 0xC8, 0xE8), wrap=True)

    txbox(s, Inches(0.6), Inches(3.3), Inches(11.5), Inches(1.1),
          text='11 releases  ·  9 modules  ·  3 architecture migrations  ·  31 calendar days  ·  105 mins measured interaction time.\n'
               'Session 3: 35 mins → 6 releases → ~60–120x compression ratio on human direction time.',
          size=Pt(11.5), color=RGBColor(0xCC, 0xDD, 0xEE), wrap=True)

    # 4 gain bullets
    gains = ['Zero handoff latency', 'Persistent cross-session context',
             'On-demand full-stack competency', 'Documentation as a side-effect']
    for i, g in enumerate(gains):
        rect(s, Inches(0.6) + i * Inches(3.0), Inches(4.4),
             Inches(2.7), Inches(0.54), fill_color=MID_BLUE)
        txbox(s, Inches(0.7) + i * Inches(3.0), Inches(4.45),
              Inches(2.5), Inches(0.46),
              text=g, size=Pt(9.5), color=WHITE, wrap=True)

    txbox(s, Inches(0.6), Inches(5.2), Inches(11.5), Inches(0.6),
          text='The gains are structural — not simply "Claude writes code faster."',
          size=Pt(13), italic=True, color=TEAL)

    # Footer
    rect(s, 0, H - Inches(0.3), W, Inches(0.3), fill_color=RGBColor(0x00, 0x1A, 0x3A))
    txbox(s, Inches(0.3), H - Inches(0.28), W - Inches(0.6), Inches(0.26),
          text='© 2026 sunder-vasudevan. All rights reserved.',
          size=Pt(7), color=RGBColor(0x88, 0x99, 0xAA))


def slide_10_recommendations(prs):
    """Slide 10 — Recommendations / Q&A."""
    s = blank_slide(prs)
    section_label(s, '09  Recommendations')
    slide_title(s, 'How to Apply This Model to Your Team')
    rule_line(s, Inches(1.2))

    recs = [
        ('1', 'Invest in context management',
         'A live NOTES.md session-start document eliminates the 15–30 min catch-up '
         'cost at the start of every AI session. This is the single highest-leverage practice.'),
        ('2', 'Treat Claude as a full-session collaborator',
         'Piecemeal queries do not capture the compounding gains. Sustained sessions '
         'where Claude holds architectural context are where the multiplier lives.'),
        ('3', 'Keep the Product Owner in the architecture seat',
         'Define what to build and why before asking Claude to help build it. '
         'Product direction and quality gate remain the human\'s responsibility.'),
        ('4', 'Log interaction time from day one',
         'This project only began logging at v4.7.0. A complete log from the start would '
         'make the efficiency analysis definitive. Use a simple markdown table per session.'),
    ]

    for i, (num, title, desc) in enumerate(recs):
        col = i % 2
        row = i // 2
        x = Inches(0.5) + col * Inches(6.3)
        y = Inches(1.45) + row * Inches(2.55)
        rect(s, x, y, Inches(5.9), Inches(2.3),
             fill_color=LIGHT_GREY if row == 0 else RGBColor(0xE8, 0xF2, 0xFC))
        txbox(s, x + Inches(0.15), y + Inches(0.12), Inches(0.45), Inches(0.45),
              text=num, size=Pt(18), bold=True, color=MID_BLUE)
        txbox(s, x + Inches(0.7), y + Inches(0.12), Inches(5.0), Inches(0.38),
              text=title, size=Pt(11), bold=True, color=DARK_NAVY)
        txbox(s, x + Inches(0.15), y + Inches(0.6), Inches(5.6), Inches(1.55),
              text=desc, size=Pt(9.5), color=DARK_TEXT, wrap=True)

    footer_bar(s)


# ── Main ──────────────────────────────────────────────────────────────────────

def build_deck():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_01_cover(prs)
    slide_02_context(prs)
    slide_03_experiment(prs)
    slide_04_what_was_built(prs)
    slide_05_velocity(prs)
    slide_05b_measured_time(prs)
    slide_06_comparison(prs)
    slide_07_how_it_works(prs)
    slide_08_quality(prs)
    slide_09_conclusion(prs)
    slide_10_recommendations(prs)

    out = '/Users/scottvalentino/BzHub/documentation/BzHub_Efficiency_ExecDeck.pptx'
    prs.save(out)
    print(f'✓  Saved: {out}  ({len(prs.slides)} slides)')


if __name__ == '__main__':
    build_deck()
